"""
Generate PPT: Perancangan Monitoring & Management Sistem — NutriGrow
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ═══════════════════════════════════════════════════════════
# COLOR PALETTE
# ═══════════════════════════════════════════════════════════
BG_DARK       = RGBColor(0x0F, 0x17, 0x2A)  # deep navy
BG_CARD       = RGBColor(0x16, 0x20, 0x3A)  # card surface
ACCENT_GREEN  = RGBColor(0x34, 0xD3, 0x99)  # emerald
ACCENT_CYAN   = RGBColor(0x22, 0xD3, 0xEE)  # cyan
ACCENT_BLUE   = RGBColor(0x60, 0xA5, 0xFA)  # sky blue
ACCENT_PURPLE = RGBColor(0xA7, 0x8B, 0xFA)  # lavender
ACCENT_ORANGE = RGBColor(0xFB, 0xBF, 0x24)  # amber
ACCENT_RED    = RGBColor(0xF8, 0x71, 0x71)  # soft red
TEXT_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_MUTED    = RGBColor(0x94, 0xA3, 0xB8)
TEXT_DIM      = RGBColor(0x64, 0x74, 0x8B)
DIVIDER       = RGBColor(0x1E, 0x29, 0x3B)

# ═══════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════
def add_bg(slide, color=BG_DARK):
    """Fill entire slide with a solid color background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None, corner_radius=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if corner_radius is not None:
        shape.adjustments[0] = corner_radius
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Segoe UI'):
    """Add a text box with a single run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox, tf

def add_multiline_text(slide, left, top, width, height, lines, default_size=14, default_color=TEXT_WHITE, font_name='Segoe UI'):
    """Add a text box with multiple paragraphs (lines is list of dicts: text, size, color, bold, alignment)."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = line_data.get('alignment', PP_ALIGN.LEFT)
        p.space_before = Pt(line_data.get('space_before', 0))
        p.space_after = Pt(line_data.get('space_after', 0))
        run = p.add_run()
        run.text = line_data.get('text', '')
        run.font.size = Pt(line_data.get('size', default_size))
        run.font.color.rgb = line_data.get('color', default_color)
        run.font.bold = line_data.get('bold', False)
        run.font.name = line_data.get('font', font_name)
    return txBox, tf

def add_badge(slide, left, top, text, bg_color, text_color=TEXT_WHITE, font_size=10, width=None):
    """Small colored badge/pill."""
    w = width or Inches(1.6)
    shape = add_rect(slide, left, top, w, Inches(0.35), bg_color, corner_radius=0.5)
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color
    run.font.bold = True
    run.font.name = 'Segoe UI'
    return shape

def add_divider_line(slide, left, top, width, color=DIVIDER):
    """Horizontal thin line."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_icon_card(slide, left, top, width, height, icon_text, title, desc, accent_color, desc_color=TEXT_MUTED):
    """A glass-style info card with icon, title, description."""
    card = add_rect(slide, left, top, width, height, BG_CARD, border_color=DIVIDER, corner_radius=0.08)
    # Icon circle
    icon_size = Inches(0.55)
    icon_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.25), top + Inches(0.25), icon_size, icon_size
    )
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = accent_color
    icon_shape.line.fill.background()
    icon_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = icon_shape.text_frame.paragraphs[0].add_run()
    r.text = icon_text
    r.font.size = Pt(16)
    r.font.color.rgb = TEXT_WHITE
    r.font.name = 'Segoe UI Emoji'

    # Title
    add_text_box(slide, left + Inches(0.25), top + Inches(0.95), width - Inches(0.5), Inches(0.35),
                 title, font_size=13, color=TEXT_WHITE, bold=True)
    # Description
    add_text_box(slide, left + Inches(0.25), top + Inches(1.3), width - Inches(0.5), height - Inches(1.5),
                 desc, font_size=10, color=desc_color)
    return card


# ═══════════════════════════════════════════════════════════
# SLIDE 1: COVER
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Decorative accent bar at top
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_GREEN)

# Small decorative circles
for (x, y, sz, c) in [
    (Inches(1), Inches(1.2), Inches(0.8), RGBColor(0x34, 0xD3, 0x99)),
    (Inches(11.5), Inches(5.8), Inches(1.0), RGBColor(0x22, 0xD3, 0xEE)),
    (Inches(0.5), Inches(6.0), Inches(0.5), RGBColor(0xA7, 0x8B, 0xFA)),
]:
    o = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, sz, sz)
    o.fill.solid()
    o.fill.fore_color.rgb = c
    # Make very transparent by setting alpha
    o.line.fill.background()
    o.fill.fore_color.brightness = 0.7

# Main Title
add_multiline_text(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(2.5), [
    {'text': 'Perancangan Monitoring &', 'size': 40, 'color': TEXT_WHITE, 'bold': True, 'space_after': 2},
    {'text': 'Management Sistem', 'size': 40, 'color': ACCENT_GREEN, 'bold': True, 'space_after': 12},
    {'text': 'Implementasi pada Proyek NutriGrow — Sistem IoT Pertanian Cerdas', 'size': 18, 'color': TEXT_MUTED, 'space_after': 8},
])

add_divider_line(slide, Inches(1.5), Inches(4.5), Inches(4))

add_multiline_text(slide, Inches(1.5), Inches(4.7), Inches(8), Inches(1.8), [
    {'text': 'Mata Kuliah: Perancangan Monitoring & Management Sistem', 'size': 14, 'color': TEXT_MUTED, 'space_after': 4},
    {'text': 'Disusun oleh: Pandya Andy Marcellino', 'size': 14, 'color': TEXT_WHITE, 'bold': True, 'space_after': 2},
    {'text': '2026', 'size': 12, 'color': TEXT_DIM, 'space_after': 2},
])

# Tech badges
badge_data = [
    ('Next.js', ACCENT_BLUE), ('Supabase', ACCENT_GREEN), ('MQTT', ACCENT_PURPLE),
    ('ESP32', ACCENT_ORANGE), ('Fuzzy Logic', ACCENT_CYAN)
]
for i, (label, color) in enumerate(badge_data):
    add_badge(slide, Inches(1.5) + Inches(i * 1.8), Inches(6.3), label, color, font_size=9, width=Inches(1.5))


# ═══════════════════════════════════════════════════════════
# SLIDE 2: OUTLINE / DAFTAR ISI
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_CYAN)

add_text_box(slide, Inches(1), Inches(0.6), Inches(10), Inches(0.6),
             'Daftar Isi', font_size=30, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.15), Inches(10), Inches(0.4),
             'Struktur presentasi mengacu pada ketentuan tugas', font_size=14, color=TEXT_MUTED)

outline_items = [
    ('01', 'Penentuan Jenis Management', 'Fault, Configuration, Security Management (ISO)', ACCENT_GREEN),
    ('02', 'Penentuan Parameter Monitoring', 'Device Uptime, RSSI, Battery Level, Actuator Duration', ACCENT_CYAN),
    ('03', 'Desain Sistem Monitoring', 'Arsitektur custom full-stack: Next.js + Supabase + MQTT', ACCENT_BLUE),
    ('04', 'Simulasi Gangguan Jaringan', '2 Skenario: Sensor Offline & Gateway Down', ACCENT_ORANGE),
    ('05', 'Dashboard Monitoring', 'Status device, grafik sensor, alert, dan trafik log', ACCENT_PURPLE),
    ('06', 'Kesimpulan & Demo', 'Ringkasan dan demonstrasi live', ACCENT_RED),
]

for i, (num, title, desc, color) in enumerate(outline_items):
    y = Inches(1.9) + Inches(i * 0.85)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), y, Inches(0.55), Inches(0.55))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    circle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = circle.text_frame.paragraphs[0].add_run()
    r.text = num
    r.font.size = Pt(16)
    r.font.bold = True
    r.font.color.rgb = TEXT_WHITE
    r.font.name = 'Segoe UI'

    add_text_box(slide, Inches(1.8), y - Inches(0.02), Inches(8), Inches(0.35),
                 title, font_size=16, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(1.8), y + Inches(0.32), Inches(8), Inches(0.3),
                 desc, font_size=11, color=TEXT_MUTED)


# ═══════════════════════════════════════════════════════════
# SLIDE 3: Penentuan Jenis Management
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_GREEN)

add_badge(slide, Inches(1), Inches(0.5), '01  JENIS MANAGEMENT', ACCENT_GREEN, font_size=9, width=Inches(2.5))
add_text_box(slide, Inches(1), Inches(1.05), Inches(10), Inches(0.6),
             'Penentuan Jenis Management', font_size=28, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.4),
             'Mengacu pada konsep ISO Network Management Framework', font_size=14, color=TEXT_MUTED)

# Three management cards
cards = [
    ('⚠️', 'Fault Management', 
     'Heartbeat monitoring via MQTT.\nJika device tidak mengirim sinyal dalam 1 menit, status otomatis berubah menjadi Offline.\nNotifikasi alert dikirim ke dashboard admin.\nPredictive maintenance via Battery Level.', ACCENT_RED),
    ('⚙️', 'Configuration Management', 
     'Pendataan Device berdasarkan ID unik & tipe.\nPemetaan perangkat ke zona/lahan.\nManajemen firmware version.\nKonfigurasi resep nutrisi & jadwal penyiraman (Fuzzy Logic + Smart Delay).', ACCENT_BLUE),
    ('🔒', 'Security Management', 
     'Autentikasi Google OAuth.\nRole-Based Access Control (RBAC):\n• Super Admin\n• Pemilik Kebun\n• Guest\nPending Approval untuk user baru.\nAudit trail pada setiap Override.', ACCENT_PURPLE),
]

card_w = Inches(3.5)
card_h = Inches(3.8)
gap = Inches(0.4)
start_x = Inches(1)

for i, (icon, title, desc, color) in enumerate(cards):
    x = start_x + (card_w + gap) * i
    add_icon_card(slide, x, Inches(2.4), card_w, card_h, icon, title, desc, color)


# ═══════════════════════════════════════════════════════════
# SLIDE 4: Penentuan Parameter Monitoring
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_CYAN)

add_badge(slide, Inches(1), Inches(0.5), '02  PARAMETER MONITORING', ACCENT_CYAN, font_size=9, width=Inches(2.8))
add_text_box(slide, Inches(1), Inches(1.05), Inches(10), Inches(0.6),
             'Penentuan Parameter Monitoring', font_size=28, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.4),
             '4 parameter kunci yang dipantau secara real-time', font_size=14, color=TEXT_MUTED)

params = [
    ('📊', 'Device Uptime & Status', 'Fault Management',
     'Mengukur keandalan dan stabilitas Sensor Node maupun Gateway. Jika gateway mati, penjadwalan penyiraman cerdas tidak akan berjalan.', ACCENT_GREEN),
    ('📶', 'RSSI (Signal Quality)', 'Performance Management',
     'Memantau kekuatan sinyal nirkabel (WiFi/ESP-NOW) di area perkebunan. Sinyal lemah menyebabkan delay atau kegagalan perintah aktuator.', ACCENT_CYAN),
    ('🔋', 'Battery Level', 'Fault Management',
     'Sensor Node nirkabel di lahan pertanian. Memonitor baterai agar pengelola bisa melakukan charging sebelum alat mati total.', ACCENT_ORANGE),
    ('⏱️', 'Actuator Duration', 'Accounting Management',
     'Memantau durasi operasi Pompa Air, Pompa Pupuk, dan Solenoid. Berguna untuk mengukur efisiensi sistem elektrik.', ACCENT_BLUE),
]

card_w = Inches(2.7)
card_h = Inches(3.8)
gap = Inches(0.3)
start_x = Inches(0.7)

for i, (icon, title, category, desc, color) in enumerate(params):
    x = start_x + (card_w + gap) * i
    y = Inches(2.3)
    
    card = add_rect(slide, x, y, card_w, card_h, BG_CARD, border_color=DIVIDER, corner_radius=0.08)
    
    # Icon
    icon_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.25), Inches(0.55), Inches(0.55))
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = color
    icon_shape.line.fill.background()
    icon_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = icon_shape.text_frame.paragraphs[0].add_run()
    r.text = icon
    r.font.size = Pt(16)
    r.font.name = 'Segoe UI Emoji'
    
    # Title
    add_text_box(slide, x + Inches(0.2), y + Inches(0.95), card_w - Inches(0.4), Inches(0.4),
                 title, font_size=13, color=TEXT_WHITE, bold=True)
    # Category badge
    add_badge(slide, x + Inches(0.2), y + Inches(1.4), category, color, font_size=8, width=Inches(2.2))
    # Description
    add_text_box(slide, x + Inches(0.2), y + Inches(1.9), card_w - Inches(0.4), Inches(1.6),
                 desc, font_size=10, color=TEXT_MUTED)


# ═══════════════════════════════════════════════════════════
# SLIDE 5: Desain Sistem Monitoring
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_BLUE)

add_badge(slide, Inches(1), Inches(0.5), '03  DESAIN SISTEM', ACCENT_BLUE, font_size=9, width=Inches(2.2))
add_text_box(slide, Inches(1), Inches(1.05), Inches(10), Inches(0.6),
             'Desain Sistem Monitoring', font_size=28, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.4),
             'Custom-built full-stack — Tanpa tools pihak ketiga (Zabbix/Grafana/PRTG)', font_size=14, color=TEXT_MUTED)

# Architecture flow cards
arch_layers = [
    ('🌱', 'Hardware Layer', 'ESP32 Sensor Node\nESP8266 Gateway\nRelay, Solenoid, Pompa', ACCENT_GREEN),
    ('📡', 'Protokol Komunikasi', 'MQTT (Mosquitto Broker)\nHTTP REST API\nWebSocket Realtime', ACCENT_CYAN),
    ('🗄️', 'Backend & Database', 'Supabase (PostgreSQL)\nRealtime subscriptions\nRow Level Security', ACCENT_PURPLE),
    ('🖥️', 'Frontend Dashboard', 'Next.js 16 (React)\nTailwind CSS\nResponsive & PWA-ready', ACCENT_BLUE),
]

card_w = Inches(2.7)
card_h = Inches(3.0)
gap = Inches(0.3)
start_x = Inches(0.7)
y_arch = Inches(2.4)

for i, (icon, title, desc, color) in enumerate(arch_layers):
    x = start_x + (card_w + gap) * i
    add_icon_card(slide, x, y_arch, card_w, card_h, icon, title, desc, color)
    # Arrow between cards
    if i < len(arch_layers) - 1:
        arrow_x = x + card_w + Inches(0.02)
        add_text_box(slide, arrow_x, y_arch + Inches(1.0), Inches(0.25), Inches(0.4),
                     '→', font_size=24, color=TEXT_DIM, alignment=PP_ALIGN.CENTER)

# Bottom note
add_rect(slide, Inches(0.7), Inches(5.8), Inches(11.8), Inches(1.0), BG_CARD, border_color=DIVIDER, corner_radius=0.06)
add_text_box(slide, Inches(1.1), Inches(5.95), Inches(11), Inches(0.7),
             '💡 Alasan tidak menggunakan Zabbix/Grafana: Sistem NutriGrow membutuhkan integrasi langsung antara logika fuzzy penyiraman, kontrol aktuator, dan monitoring sensor dalam satu platform terpadu. Tools pihak ketiga hanya menyediakan visualisasi, bukan pengendalian end-to-end.',
             font_size=11, color=TEXT_MUTED)


# ═══════════════════════════════════════════════════════════
# SLIDE 6: Simulasi Gangguan Jaringan
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_ORANGE)

add_badge(slide, Inches(1), Inches(0.5), '04  SIMULASI GANGGUAN', ACCENT_ORANGE, font_size=9, width=Inches(2.5))
add_text_box(slide, Inches(1), Inches(1.05), Inches(10), Inches(0.6),
             'Simulasi Gangguan Jaringan', font_size=28, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.4),
             'Minimal 2 skenario gangguan untuk demonstrasi live', font_size=14, color=TEXT_MUTED)

# Scenario 1
sc1_x, sc1_y = Inches(0.7), Inches(2.3)
sc1_w, sc1_h = Inches(5.8), Inches(4.5)
add_rect(slide, sc1_x, sc1_y, sc1_w, sc1_h, BG_CARD, border_color=ACCENT_RED, corner_radius=0.06)
add_badge(slide, sc1_x + Inches(0.3), sc1_y + Inches(0.3), 'SKENARIO 1', ACCENT_RED, font_size=9, width=Inches(1.5))
add_text_box(slide, sc1_x + Inches(0.3), sc1_y + Inches(0.75), Inches(5), Inches(0.4),
             '⚡ Sensor Node Offline', font_size=18, color=TEXT_WHITE, bold=True)
add_multiline_text(slide, sc1_x + Inches(0.3), sc1_y + Inches(1.3), Inches(5.2), Inches(3.0), [
    {'text': 'Cara Simulasi:', 'size': 12, 'color': ACCENT_GREEN, 'bold': True, 'space_after': 2},
    {'text': 'Cabut power / matikan ESP32 Sensor Node.', 'size': 11, 'color': TEXT_MUTED, 'space_after': 8},
    {'text': 'Deteksi Sistem:', 'size': 12, 'color': ACCENT_GREEN, 'bold': True, 'space_after': 2},
    {'text': 'Server mendeteksi heartbeat terputus. Setelah 1 menit tanpa sinyal, status otomatis berubah menjadi Offline.', 'size': 11, 'color': TEXT_MUTED, 'space_after': 8},
    {'text': 'Alert ke Admin:', 'size': 12, 'color': ACCENT_GREEN, 'bold': True, 'space_after': 2},
    {'text': '• Indikator di Dashboard berubah merah (Offline)\n• Notifikasi device_alert muncul di menu Notifications\n• Semua sensor (Soil Moisture, DHT22, TDS, pH) di zona tersebut ikut Offline', 'size': 11, 'color': TEXT_MUTED},
])

# Scenario 2
sc2_x = Inches(6.8)
sc2_w = Inches(5.8)
add_rect(slide, sc2_x, sc1_y, sc2_w, sc1_h, BG_CARD, border_color=ACCENT_ORANGE, corner_radius=0.06)
add_badge(slide, sc2_x + Inches(0.3), sc1_y + Inches(0.3), 'SKENARIO 2', ACCENT_ORANGE, font_size=9, width=Inches(1.5))
add_text_box(slide, sc2_x + Inches(0.3), sc1_y + Inches(0.75), Inches(5), Inches(0.4),
             '🌐 Gateway Down', font_size=18, color=TEXT_WHITE, bold=True)
add_multiline_text(slide, sc2_x + Inches(0.3), sc1_y + Inches(1.3), Inches(5.2), Inches(3.0), [
    {'text': 'Cara Simulasi:', 'size': 12, 'color': ACCENT_CYAN, 'bold': True, 'space_after': 2},
    {'text': 'Matikan router WiFi atau matikan ESP8266 Gateway.', 'size': 11, 'color': TEXT_MUTED, 'space_after': 8},
    {'text': 'Deteksi Sistem:', 'size': 12, 'color': ACCENT_CYAN, 'bold': True, 'space_after': 2},
    {'text': 'Koneksi MQTT Broker terputus dari sisi device. Semua sensor kehilangan jembatan data.', 'size': 11, 'color': TEXT_MUTED, 'space_after': 8},
    {'text': 'Dampak Lanjutan:', 'size': 12, 'color': ACCENT_CYAN, 'bold': True, 'space_after': 2},
    {'text': '• Status Gateway di Dashboard otomatis berubah Offline\n• Semua Sensor Node ikut Offline (efek domino)\n• Perintah Override dari web tidak sampai ke hardware\n• Web tetap bisa diakses (data historis tetap ada)', 'size': 11, 'color': TEXT_MUTED},
])


# ═══════════════════════════════════════════════════════════
# SLIDE 7: Dashboard Monitoring
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_PURPLE)

add_badge(slide, Inches(1), Inches(0.5), '05  DASHBOARD MONITORING', ACCENT_PURPLE, font_size=9, width=Inches(2.8))
add_text_box(slide, Inches(1), Inches(1.05), Inches(10), Inches(0.6),
             'Dashboard Monitoring', font_size=28, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.4),
             'Semua opsi yang diminta telah terimplementasi', font_size=14, color=TEXT_MUTED)

features = [
    ('📱', 'Status Device\n& Online/Offline', 'Indikator warna real-time per perangkat individual:\nSoil Moisture, DHT22, TDS, pH, Relay, Solenoid, Pompa Air, Pompa Pupuk.\nBadge Online (hijau) / Offline (merah) dengan sisa baterai.', ACCENT_GREEN),
    ('📊', 'Grafik Sensor\nReal-time', 'Visualisasi data telemetry:\nSuhu, Kelembaban Tanah, Kelembaban Udara, pH, TDS.\nGrafik historis 24 jam dengan tren.\nUpdate otomatis via WebSocket subscription.', ACCENT_CYAN),
    ('📝', 'Trafik & Log\nAktivitas', 'Mencatat kapan aktuator dinyalakan.\nDurasi operasi setiap pompa & solenoid.\nLog Override dengan nama operator.\nRiwayat irigasi otomatis & manual.', ACCENT_BLUE),
    ('🔔', 'Alert &\nNotifikasi', 'Lonceng notifikasi di header Dashboard.\nRiwayat anomali (device offline alert).\nPutusnya koneksi perangkat tercatat.\nBadge count untuk notifikasi belum dibaca.', ACCENT_ORANGE),
]

card_w = Inches(2.7)
card_h = Inches(3.8)
gap = Inches(0.3)
start_x = Inches(0.7)

for i, (icon, title, desc, color) in enumerate(features):
    x = start_x + (card_w + gap) * i
    add_icon_card(slide, x, Inches(2.3), card_w, card_h, icon, title, desc, color)

# Checkmark row at bottom
add_rect(slide, Inches(0.7), Inches(6.5), Inches(11.8), Inches(0.55), BG_CARD, border_color=ACCENT_GREEN, corner_radius=0.06)
add_text_box(slide, Inches(1.1), Inches(6.55), Inches(11), Inches(0.45),
             '✅ Status Device    ✅ Online/Offline    ✅ Trafik & Log    ✅ Grafik Sensor    ✅ Alert — Semua opsi dari ketentuan tugas terpenuhi',
             font_size=12, color=ACCENT_GREEN, bold=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 8: Device Management Detail
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_BLUE)

add_badge(slide, Inches(1), Inches(0.5), '05b  DEVICE MANAGEMENT', ACCENT_BLUE, font_size=9, width=Inches(2.8))
add_text_box(slide, Inches(1), Inches(1.05), Inches(10), Inches(0.6),
             'Perangkat Individual per Zona', font_size=28, color=TEXT_WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(1.65), Inches(10), Inches(0.4),
             'Setiap sensor dan aktuator memiliki kartu status masing-masing', font_size=14, color=TEXT_MUTED)

# Sensors section
add_text_box(slide, Inches(0.7), Inches(2.25), Inches(5), Inches(0.4),
             '📡 Sensor Nodes', font_size=18, color=ACCENT_GREEN, bold=True)

sensor_devices = [
    ('💧', 'Soil Moisture', 'Kelembaban tanah kapasitif', RGBColor(0x22, 0xD3, 0xEE)),
    ('🌡️', 'DHT22', 'Suhu & kelembaban udara', RGBColor(0xFB, 0x92, 0x3C)),
    ('🧪', 'TDS', 'Total Dissolved Solids / EC', RGBColor(0x2D, 0xD4, 0xBF)),
    ('⚗️', 'pH', 'Kadar asam-basa nutrisi', RGBColor(0x81, 0x8C, 0xF8)),
]

for i, (icon, name, desc, color) in enumerate(sensor_devices):
    y = Inches(2.8) + Inches(i * 0.95)
    card = add_rect(slide, Inches(0.7), y, Inches(5.5), Inches(0.8), BG_CARD, border_color=DIVIDER, corner_radius=0.06)
    
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.95), y + Inches(0.12), Inches(0.55), Inches(0.55))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = circ.text_frame.paragraphs[0].add_run()
    r.text = icon
    r.font.size = Pt(16)
    r.font.name = 'Segoe UI Emoji'
    
    add_text_box(slide, Inches(1.7), y + Inches(0.08), Inches(3.5), Inches(0.35),
                 name, font_size=14, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(1.7), y + Inches(0.42), Inches(3.5), Inches(0.3),
                 desc, font_size=10, color=TEXT_MUTED)

# Actuators section
add_text_box(slide, Inches(6.8), Inches(2.25), Inches(5), Inches(0.4),
             '⚙️ Aktuator Nodes', font_size=18, color=ACCENT_ORANGE, bold=True)

actuator_devices = [
    ('🔌', 'Relay', 'Kontrol switch utama', RGBColor(0xFB, 0xBF, 0x24)),
    ('🚰', 'Solenoid Valve', 'Katup distribusi air', RGBColor(0x60, 0xA5, 0xFA)),
    ('💦', 'Pompa Air', 'Pompa distribusi air bersih', RGBColor(0x38, 0xBD, 0xF8)),
    ('🧴', 'Pompa Pupuk', 'Pompa larutan nutrisi/fertigasi', RGBColor(0xA3, 0xE6, 0x35)),
]

for i, (icon, name, desc, color) in enumerate(actuator_devices):
    y = Inches(2.8) + Inches(i * 0.95)
    card = add_rect(slide, Inches(6.8), y, Inches(5.5), Inches(0.8), BG_CARD, border_color=DIVIDER, corner_radius=0.06)
    
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.05), y + Inches(0.12), Inches(0.55), Inches(0.55))
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    circ.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = circ.text_frame.paragraphs[0].add_run()
    r.text = icon
    r.font.size = Pt(16)
    r.font.name = 'Segoe UI Emoji'
    
    add_text_box(slide, Inches(7.8), y + Inches(0.08), Inches(3.5), Inches(0.35),
                 name, font_size=14, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(7.8), y + Inches(0.42), Inches(3.5), Inches(0.3),
                 desc, font_size=10, color=TEXT_MUTED)

# Gateway at bottom
add_rect(slide, Inches(0.7), Inches(6.7), Inches(11.8), Inches(0.55), BG_CARD, border_color=ACCENT_PURPLE, corner_radius=0.06)
add_text_box(slide, Inches(1.1), Inches(6.75), Inches(11), Inches(0.45),
             '📻 ESP8266 Gateway — Status dinamis: Online jika minimal 1 sensor aktif, Offline jika semua sensor mati (efek domino)',
             font_size=12, color=ACCENT_PURPLE, bold=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 9: Kesimpulan
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_GREEN)

add_badge(slide, Inches(1), Inches(0.5), '06  KESIMPULAN', ACCENT_GREEN, font_size=9, width=Inches(2))
add_text_box(slide, Inches(1), Inches(1.05), Inches(10), Inches(0.6),
             'Kesimpulan & Demo', font_size=28, color=TEXT_WHITE, bold=True)

# Summary table
summary_items = [
    ('01', 'Jenis Management', '3 pilar ISO: Fault, Configuration, Security', '✅', ACCENT_GREEN),
    ('02', 'Parameter Monitoring', '4 parameter: Uptime, RSSI, Battery, Duration', '✅', ACCENT_CYAN),
    ('03', 'Desain Sistem', 'Custom full-stack (Next.js + Supabase + MQTT)', '✅', ACCENT_BLUE),
    ('04', 'Simulasi Gangguan', '2 skenario: Sensor Offline & Gateway Down', '✅', ACCENT_ORANGE),
    ('05', 'Dashboard Monitoring', 'Status, Grafik, Alert, Trafik — Semua terpenuhi', '✅', ACCENT_PURPLE),
]

for i, (num, title, desc, status, color) in enumerate(summary_items):
    y = Inches(1.9) + Inches(i * 0.85)
    row_card = add_rect(slide, Inches(1), y, Inches(11), Inches(0.7), BG_CARD, border_color=DIVIDER, corner_radius=0.06)
    
    # Number
    add_text_box(slide, Inches(1.3), y + Inches(0.12), Inches(0.5), Inches(0.4),
                 num, font_size=16, color=color, bold=True)
    # Title
    add_text_box(slide, Inches(2), y + Inches(0.08), Inches(3), Inches(0.35),
                 title, font_size=14, color=TEXT_WHITE, bold=True)
    # Description
    add_text_box(slide, Inches(2), y + Inches(0.38), Inches(7), Inches(0.3),
                 desc, font_size=11, color=TEXT_MUTED)
    # Status check
    add_text_box(slide, Inches(10.8), y + Inches(0.12), Inches(0.8), Inches(0.4),
                 status, font_size=20, color=ACCENT_GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# Demo prompt
add_rect(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.8), BG_CARD, border_color=ACCENT_GREEN, corner_radius=0.06)
add_text_box(slide, Inches(1.4), Inches(6.4), Inches(10), Inches(0.6),
             '🎬  DEMO LIVE — Sistem siap didemonstrasikan secara langsung. Silakan siapkan hardware (ESP32 + ESP8266) untuk menunjukkan skenario gangguan secara real-time.',
             font_size=13, color=ACCENT_GREEN, bold=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 10: TERIMA KASIH
# ═══════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.06), ACCENT_GREEN)

add_multiline_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(3.0), [
    {'text': 'Terima Kasih', 'size': 44, 'color': TEXT_WHITE, 'bold': True, 'alignment': PP_ALIGN.CENTER, 'space_after': 8},
    {'text': '— Sesi Tanya Jawab —', 'size': 20, 'color': ACCENT_GREEN, 'alignment': PP_ALIGN.CENTER, 'space_after': 20},
    {'text': 'NutriGrow — Sistem IoT Pertanian Cerdas', 'size': 16, 'color': TEXT_MUTED, 'alignment': PP_ALIGN.CENTER, 'space_after': 4},
    {'text': 'Pandya Andy Marcellino', 'size': 14, 'color': TEXT_DIM, 'alignment': PP_ALIGN.CENTER},
])

# Tech badges at bottom center
badge_data2 = [
    ('Next.js 16', ACCENT_BLUE), ('Supabase', ACCENT_GREEN), ('MQTT', ACCENT_PURPLE),
    ('ESP32 + ESP8266', ACCENT_ORANGE), ('Fuzzy Logic', ACCENT_CYAN)
]
total_w = len(badge_data2) * Inches(1.8)
start_badge = (prs.slide_width - total_w) / 2
for i, (label, color) in enumerate(badge_data2):
    bx = start_badge + Inches(i * 1.8)
    add_badge(slide, int(bx), Inches(5.5), label, color, font_size=9, width=Inches(1.6))


# ═══════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════
output_path = r'c:\Users\Pandya\Downloads\nutrigrow-web\nutrigrow-web\Laporan_Monitoring_Management_NutriGrow.pptx'
prs.save(output_path)
print(f'PPT berhasil disimpan: {output_path}')
print(f'Total slide: {len(prs.slides)}')
